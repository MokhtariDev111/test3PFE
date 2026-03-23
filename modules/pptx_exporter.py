"""
pptx_exporter.py  — v2.1: Premium Templates (corruption-free)
==============================================================
Fixes vs v2:
  1. _add_transition: use .append() not .insert(-1) — insert at wrong
     position breaks XML schema order PowerPoint enforces.
  2. Progress bar: guard fill_w > 10 before drawing fill rect —
     a zero/tiny width shape produces invalid EMU values in XML.
  3. Removed negative-x decorative circle — shapes with x < 0 are
     outside the slide canvas and trigger PowerPoint repair warnings.
  4. Notes: wrapped in try/except with has_notes_slide check.
  5. BytesIO streams: seek(0) before passing to add_picture() to
     prevent empty image elements after stream already consumed.
  6. _add_rect guard: skip if width or height is zero.
"""

import logging, re, io, base64
from dataclasses import dataclass, field
from pathlib import Path
from datetime import datetime
import sys

ROOT_DIR = Path(__file__).resolve().parent.parent
if str(ROOT_DIR) not in sys.path:
    sys.path.insert(0, str(ROOT_DIR))

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

try:
    from modules.slide_generator import SlideData
except ImportError:
    pass

from modules.config_loader import CONFIG

log = logging.getLogger("pptx_exporter")
if not log.hasHandlers():
    logging.basicConfig(level=logging.INFO,
                        format="%(asctime)s [%(levelname)s] %(name)s — %(message)s",
                        datefmt="%H:%M:%S")


# ── Theme System ───────────────────────────────────────────────────────────────
@dataclass
class ThemeConfig:
    name:      str
    bg:        RGBColor
    bg2:       RGBColor
    accent:    RGBColor
    accent2:   RGBColor
    text:      RGBColor
    subtext:   RGBColor
    bullet:    RGBColor
    number_bg: RGBColor


def _rgb(r, g, b): return RGBColor(r, g, b)


THEMES: dict[str, ThemeConfig] = {
    "Dark Navy": ThemeConfig(
        name="Dark Navy",
        bg=_rgb(0x0D,0x1B,0x2A), bg2=_rgb(0x16,0x2A,0x40),
        accent=_rgb(0x1F,0x6F,0xEB), accent2=_rgb(0x00,0xC6,0xFF),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xAA,0xC4,0xD8),
        bullet=_rgb(0xC8,0xE6,0xFF), number_bg=_rgb(0x1F,0x6F,0xEB),
    ),
    "Crimson Pro": ThemeConfig(
        name="Crimson Pro",
        bg=_rgb(0x1A,0x07,0x07), bg2=_rgb(0x2B,0x0E,0x0E),
        accent=_rgb(0xE0,0x3A,0x2F), accent2=_rgb(0xFF,0x6B,0x6B),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xF5,0xC6,0xC6),
        bullet=_rgb(0xFF,0xCC,0xCC), number_bg=_rgb(0xE0,0x3A,0x2F),
    ),
    "Forest Green": ThemeConfig(
        name="Forest Green",
        bg=_rgb(0x07,0x18,0x0D), bg2=_rgb(0x0E,0x28,0x16),
        accent=_rgb(0x27,0xAE,0x60), accent2=_rgb(0x52,0xD9,0x68),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xC3,0xE6,0xC3),
        bullet=_rgb(0xB7,0xF5,0xC8), number_bg=_rgb(0x27,0xAE,0x60),
    ),
    "Midnight Purple": ThemeConfig(
        name="Midnight Purple",
        bg=_rgb(0x0F,0x08,0x1A), bg2=_rgb(0x1C,0x12,0x2E),
        accent=_rgb(0x9B,0x59,0xB6), accent2=_rgb(0xCC,0x99,0xFF),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xD7,0xC4,0xE8),
        bullet=_rgb(0xE8,0xD8,0xFF), number_bg=_rgb(0x9B,0x59,0xB6),
    ),
    "Corporate White": ThemeConfig(
        name="Corporate White",
        bg=_rgb(0xFA,0xFA,0xFC), bg2=_rgb(0xEE,0xEF,0xF5),
        accent=_rgb(0x2C,0x3E,0x50), accent2=_rgb(0x34,0x98,0xDB),
        text=_rgb(0x1A,0x1A,0x2E), subtext=_rgb(0x55,0x66,0x77),
        bullet=_rgb(0x2C,0x3E,0x50), number_bg=_rgb(0x2C,0x3E,0x50),
    ),
    "Ocean Breeze": ThemeConfig(
        name="Ocean Breeze",
        bg=_rgb(0x05,0x1E,0x3E), bg2=_rgb(0x0A,0x2D,0x55),
        accent=_rgb(0x00,0xB4,0xD8), accent2=_rgb(0x90,0xE0,0xEF),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xB8,0xE0,0xF0),
        bullet=_rgb(0xCA,0xF0,0xF8), number_bg=_rgb(0x00,0xB4,0xD8),
    ),
    "Sunset Gold": ThemeConfig(
        name="Sunset Gold",
        bg=_rgb(0x1C,0x10,0x07), bg2=_rgb(0x2E,0x1A,0x0A),
        accent=_rgb(0xF0,0xA5,0x00), accent2=_rgb(0xFF,0xD1,0x66),
        text=_rgb(0xFF,0xFF,0xFF), subtext=_rgb(0xF5,0xE0,0xB0),
        bullet=_rgb(0xFF,0xEC,0xCC), number_bg=_rgb(0xF0,0xA5,0x00),
    ),
    "Slate Pro": ThemeConfig(
        name="Slate Pro",
        bg=_rgb(0x1A,0x1D,0x23), bg2=_rgb(0x25,0x29,0x33),
        accent=_rgb(0x64,0x8F,0xFF), accent2=_rgb(0xA0,0xC4,0xFF),
        text=_rgb(0xF0,0xF2,0xFF), subtext=_rgb(0xA0,0xAB,0xC4),
        bullet=_rgb(0xD0,0xDC,0xFF), number_bg=_rgb(0x64,0x8F,0xFF),
    ),
}
DEFAULT_THEME = THEMES["Dark Navy"]


# ── Low-level helpers ──────────────────────────────────────────────────────────

def _set_bg(slide, color: RGBColor):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def _add_rect(slide, left, top, width, height, color: RGBColor):
    """FIX 6: guard against zero-size shapes which corrupt XML."""
    if width <= 0 or height <= 0:
        return None
    shape = slide.shapes.add_shape(1, int(left), int(top), int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_oval(slide, left, top, width, height, color: RGBColor):
    """Safe oval — 9 = MSO oval shape type."""
    if width <= 0 or height <= 0:
        return None
    # FIX 3: clamp to slide canvas (no negative positions)
    left = max(0, int(left))
    top  = max(0, int(top))
    shape = slide.shapes.add_shape(9, left, top, int(width), int(height))
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def _add_textbox(slide, text, left, top, width, height,
                 size, bold, color: RGBColor,
                 align=PP_ALIGN.LEFT, wrap=True, italic=False):
    if not text or width <= 0 or height <= 0:
        return None
    tb = slide.shapes.add_textbox(int(left), int(top), int(width), int(height))
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = str(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return tb


def _add_transition(slide):
    """FIX 1: use .append() — transition must be last child of spTree parent."""
    from pptx.oxml import parse_xml
    xml = (
        '<p:transition xmlns:p="http://schemas.openxmlformats.org/'
        'presentationml/2006/main" spd="med"><p:fade/></p:transition>'
    )
    # append to the slide element (spTree's parent = cSld's parent = sld)
    slide.element.append(parse_xml(xml))


def _write_notes(slide_obj, text: str):
    """FIX 4: safe notes write with existence check."""
    if not text:
        return
    try:
        tf = slide_obj.notes_slide.notes_text_frame
        tf.text = text[:2000]  # cap length to avoid XML bloat
    except Exception as e:
        log.debug(f"Could not write notes: {e}")


# ── Geometric decoration ───────────────────────────────────────────────────────

def _draw_corner_accents(slide, W, H, accent: RGBColor, accent2: RGBColor):
    """Step accents top-right and bottom-left — all positions clamped > 0."""
    # Top-right steps
    _add_rect(slide, W - Inches(3.0), 0, Inches(3.0), Inches(0.07), accent)
    _add_rect(slide, W - Inches(1.6), 0, Inches(1.6), Inches(0.18), accent)
    _add_rect(slide, W - Inches(0.7), 0, Inches(0.7), Inches(0.42), accent2)
    # Bottom-left steps
    _add_rect(slide, 0, H - Inches(0.07), Inches(3.0), Inches(0.07), accent)
    _add_rect(slide, 0, H - Inches(0.18), Inches(1.6), Inches(0.18), accent)
    _add_rect(slide, 0, H - Inches(0.42), Inches(0.7), Inches(0.42), accent2)


def _draw_bg_layers(slide, W, H, bg: RGBColor, bg2: RGBColor):
    """Simulated gradient via two layered rects. FIX 3: no negative coords."""
    _add_rect(slide, 0, 0, W, H, bg)
    # Bottom-right secondary panel
    _add_rect(slide, W - Inches(3.8), H - Inches(2.2), Inches(3.8), Inches(2.2), bg2)
    # Top-left secondary panel (smaller)
    _add_rect(slide, 0, 0, Inches(2.0), Inches(1.5), bg2)
    # Decorative oval — bottom-right, fully on canvas
    _add_oval(slide, W - Inches(2.2), H - Inches(2.2), Inches(2.5), Inches(2.5), bg2)


def _draw_numbered_bullet(slide, number: int, x, y, theme: ThemeConfig):
    """Circle with number inside."""
    D = Inches(0.3)
    _add_oval(slide, int(x), int(y), int(D), int(D), theme.number_bg)
    _add_textbox(
        slide, str(number),
        x, y + Inches(0.03), D, D,
        size=9, bold=True, color=theme.text, align=PP_ALIGN.CENTER, wrap=False
    )


def _draw_progress_bar(slide, W, H, current: int, total: int,
                       accent: RGBColor, subtext: RGBColor):
    """FIX 2: guard fill_w > 0 before drawing fill rect."""
    bar_y = H - Inches(0.15)
    bar_h = Inches(0.06)
    # Track (full width)
    _add_rect(slide, 0, bar_y, W, bar_h, RGBColor(0x33,0x33,0x44))
    # Fill — only draw if nonzero width
    if total > 0:
        fill_w = int(W * current / total)
        if fill_w > 10:   # FIX 2: minimum 10 EMU to avoid zero/tiny corruption
            _add_rect(slide, 0, bar_y, fill_w, bar_h, accent)
    # Counter
    _add_textbox(
        slide, f"{current} / {total}",
        W - Inches(1.0), H - Inches(0.28), Inches(0.9), Inches(0.22),
        size=8, bold=False, color=subtext, align=PP_ALIGN.RIGHT
    )


def _draw_key_message_box(slide, text: str, left, top, width, height,
                           accent: RGBColor, theme: ThemeConfig):
    """Highlighted key message box with left accent stripe."""
    _add_rect(slide, left, top, width, height, theme.bg2)
    _add_rect(slide, left, top, Inches(0.07), height, accent)
    _add_textbox(
        slide, "KEY INSIGHT",
        left + Inches(0.15), top + Inches(0.08),
        width - Inches(0.2), Inches(0.22),
        size=8, bold=True, color=accent
    )
    _add_textbox(
        slide, f'"{text}"',
        left + Inches(0.15), top + Inches(0.28),
        width - Inches(0.25), height - Inches(0.32),
        size=13, bold=False, italic=True, color=theme.text
    )


# ── Title Slide ────────────────────────────────────────────────────────────────

def build_title_slide(prs, slide_data, theme: ThemeConfig, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_transition(slide)
    W, H = prs.slide_width, prs.slide_height

    _draw_bg_layers(slide, W, H, theme.bg, theme.bg2)
    _draw_corner_accents(slide, W, H, theme.accent, theme.accent2)

    # Left accent band
    _add_rect(slide, 0, 0, Inches(0.22), H, theme.accent)
    _add_rect(slide, Inches(0.22), 0, Inches(0.05), H, theme.accent2)

    # Top stripe
    _add_rect(slide, 0, 0, W, Inches(0.06), theme.accent2)

    # Decorative circle (top-right, fully on canvas)
    _add_oval(slide, W - Inches(3.2), Inches(0.4), Inches(2.8), Inches(2.8), theme.bg2)

    # Title
    _add_textbox(
        slide, slide_data.title,
        Inches(0.55), Inches(1.8), Inches(8.2), Inches(1.8),
        size=42, bold=True, color=theme.text
    )

    # Accent underline
    _add_rect(slide, Inches(0.55), Inches(3.62), Inches(2.2), Inches(0.05), theme.accent)

    # Date / subtitle
    _add_textbox(
        slide,
        f"AI-Generated Educational Presentation  •  {datetime.now().strftime('%B %Y')}",
        Inches(0.55), Inches(3.75), Inches(8.0), Inches(0.45),
        size=13, bold=False, color=theme.subtext
    )

    _draw_progress_bar(slide, W, H, 1, total, theme.accent, theme.subtext)
    _write_notes(slide, slide_data.speaker_notes)


# ── Section Divider Slide ──────────────────────────────────────────────────────

def build_section_slide(prs, title: str, theme: ThemeConfig,
                         current: int, total: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_transition(slide)
    W, H = prs.slide_width, prs.slide_height

    _draw_bg_layers(slide, W, H, theme.bg, theme.bg2)

    # Horizontal mid-stripe
    _add_rect(slide, 0, H//2 - Inches(0.04), W, Inches(0.07), theme.accent)
    _add_rect(slide, 0, H//2 + Inches(0.03), W, Inches(0.02), theme.accent2)

    # Large watermark number
    _add_textbox(
        slide, str(current),
        W - Inches(2.8), Inches(0.3), Inches(2.5), Inches(2.5),
        size=110, bold=True, color=theme.bg2, align=PP_ALIGN.RIGHT
    )

    # Section label
    _add_textbox(
        slide, "SECTION",
        Inches(0.7), H//2 - Inches(1.1), Inches(5.0), Inches(0.35),
        size=10, bold=True, color=theme.accent
    )
    _add_textbox(
        slide, title,
        Inches(0.7), H//2 - Inches(0.75), Inches(8.0), Inches(0.75),
        size=34, bold=True, color=theme.text
    )
    _add_textbox(
        slide, "— Overview",
        Inches(0.7), H//2 + Inches(0.08), Inches(4.0), Inches(0.35),
        size=13, bold=False, color=theme.subtext
    )

    _draw_progress_bar(slide, W, H, current, total, theme.accent, theme.subtext)


# ── Content Slides ─────────────────────────────────────────────────────────────

def build_content_slide(prs, slide_data, theme: ThemeConfig,
                         current: int, total: int,
                         diagram_path: str = None, pdf_images: dict = None):
    pdf_images = pdf_images or {}

    # Determine visual stream (FIX 5: seek(0) before use)
    visual_stream = None
    if diagram_path and Path(diagram_path).exists():
        visual_stream = diagram_path
    elif getattr(slide_data, 'image_id', None) and slide_data.image_id in pdf_images:
        img_data = base64.b64decode(pdf_images[slide_data.image_id])
        buf = io.BytesIO(img_data)
        buf.seek(0)          # FIX 5
        visual_stream = buf

    has_visual  = visual_stream is not None
    bullets     = slide_data.bullets or []
    few_bullets = len(bullets) <= 3

    if has_visual:
        _build_split_slide(prs, slide_data, theme, current, total, visual_stream, bullets)
    elif few_bullets:
        _build_twocol_slide(prs, slide_data, theme, current, total, bullets)
    else:
        _build_standard_slide(prs, slide_data, theme, current, total, bullets)


def _header(slide, title, W, theme):
    _add_rect(slide, 0, 0, W, Inches(1.05), theme.accent)
    _add_rect(slide, 0, Inches(1.05), W, Inches(0.04), theme.accent2)
    _add_textbox(
        slide, title,
        Inches(0.35), Inches(0.13), Inches(9.3), Inches(0.82),
        size=24, bold=True, color=theme.text
    )


def _build_standard_slide(prs, slide_data, theme, current, total, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_transition(slide)
    W, H = prs.slide_width, prs.slide_height

    _draw_bg_layers(slide, W, H, theme.bg, theme.bg2)
    _draw_corner_accents(slide, W, H, theme.accent, theme.accent2)
    _header(slide, slide_data.title, W, theme)

    for i, bullet in enumerate(bullets[:4]):
        txt = bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
        y   = Inches(1.22) + i * Inches(0.88)
        _draw_numbered_bullet(slide, i + 1, Inches(0.28), y + Inches(0.02), theme)
        _add_textbox(slide, txt, Inches(0.72), y, Inches(9.0), Inches(0.8),
                     size=17, bold=False, color=theme.bullet)

    km = getattr(slide_data, 'key_message', None)
    if km and len(bullets) <= 3:
        box_y = Inches(1.22) + len(bullets) * Inches(0.88) + Inches(0.1)
        _draw_key_message_box(slide, km, Inches(0.28), box_y,
                               Inches(9.4), Inches(0.85), theme.accent, theme)

    _draw_progress_bar(slide, W, H, current, total, theme.accent, theme.subtext)
    _write_notes(slide, slide_data.speaker_notes)


def _build_twocol_slide(prs, slide_data, theme, current, total, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_transition(slide)
    W, H = prs.slide_width, prs.slide_height

    _draw_bg_layers(slide, W, H, theme.bg, theme.bg2)
    _header(slide, slide_data.title, W, theme)

    # Vertical divider
    _add_rect(slide, Inches(5.55), Inches(1.12), Inches(0.04),
              H - Inches(1.32), theme.accent2)

    # Left: numbered bullets
    for i, bullet in enumerate(bullets[:4]):
        txt = bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
        y   = Inches(1.22) + i * Inches(0.88)
        _draw_numbered_bullet(slide, i + 1, Inches(0.28), y + Inches(0.02), theme)
        _add_textbox(slide, txt, Inches(0.7), y, Inches(4.65), Inches(0.82),
                     size=15, bold=False, color=theme.bullet)

    # Right: key message panel
    km = getattr(slide_data, 'key_message', None)
    if not km and bullets:
        last = bullets[-1]
        km = last.get("text", "") if isinstance(last, dict) else str(last)
    km = km or ""

    panel_x = Inches(5.72)
    panel_w = W - panel_x - Inches(0.15)
    panel_h = H - Inches(1.32)
    _add_rect(slide, panel_x, Inches(1.12), panel_w, panel_h, theme.bg2)
    _add_rect(slide, panel_x, Inches(1.12), Inches(0.07), panel_h, theme.accent)

    _add_textbox(slide, "KEY INSIGHT",
                 panel_x + Inches(0.15), Inches(1.24), panel_w - Inches(0.2), Inches(0.25),
                 size=9, bold=True, color=theme.accent)
    _add_textbox(slide, f'"{km}"' if km else "",
                 panel_x + Inches(0.15), Inches(1.55), panel_w - Inches(0.2),
                 H - Inches(2.0),
                 size=15, bold=False, italic=True, color=theme.text)

    _draw_progress_bar(slide, W, H, current, total, theme.accent, theme.subtext)
    _write_notes(slide, slide_data.speaker_notes)


def _build_split_slide(prs, slide_data, theme, current, total, visual_stream, bullets):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_transition(slide)
    W, H = prs.slide_width, prs.slide_height

    _draw_bg_layers(slide, W, H, theme.bg, theme.bg2)
    _header(slide, slide_data.title, W, theme)

    # Left: numbered bullets
    for i, bullet in enumerate(bullets[:4]):
        txt = bullet.get("text", "") if isinstance(bullet, dict) else str(bullet)
        y   = Inches(1.18) + i * Inches(0.84)
        _draw_numbered_bullet(slide, i + 1, Inches(0.28), y + Inches(0.02), theme)
        _add_textbox(slide, txt, Inches(0.7), y, Inches(5.0), Inches(0.78),
                     size=15, bold=False, color=theme.bullet)

    # Right: visual panel
    panel_x = Inches(5.82)
    panel_w = W - panel_x - Inches(0.12)
    _add_rect(slide, panel_x, Inches(1.12), panel_w, H - Inches(1.32), theme.bg2)
    try:
        # FIX 5: seek(0) for BytesIO before add_picture
        if hasattr(visual_stream, 'seek'):
            visual_stream.seek(0)
        slide.shapes.add_picture(visual_stream, panel_x + Inches(0.1),
                                  Inches(1.22), panel_w - Inches(0.2))
    except Exception as e:
        log.warning(f"Could not embed visual on slide: {e}")

    _draw_progress_bar(slide, W, H, current, total, theme.accent, theme.subtext)
    _write_notes(slide, slide_data.speaker_notes)


# ── Main export ────────────────────────────────────────────────────────────────

def export_pptx(slides, topic: str = "Presentation",
                theme: ThemeConfig = None,
                diagrams: dict = None,
                pdf_images: dict = None,
                output_dir: str = None) -> str:

    theme    = theme or DEFAULT_THEME
    diagrams = diagrams or {}

    prs = Presentation()
    prs.slide_width  = Inches(10)
    prs.slide_height = Inches(5.625)
    total = len(slides)

    for i, slide_data in enumerate(slides):
        current = i + 1
        if slide_data.slide_type == "title":
            build_title_slide(prs, slide_data, theme, total)
        elif slide_data.slide_type == "section":
            build_section_slide(prs, slide_data.title, theme, current, total)
        else:
            # Auto section divider every 3 slides on long decks
            if total > 6 and i > 0 and i % 3 == 0:
                build_section_slide(prs, slide_data.title, theme, current, total)
            build_content_slide(
                prs, slide_data, theme, current, total,
                diagrams.get(i), pdf_images
            )

    # Output path
    if output_dir:
        out_dir = Path(output_dir).resolve()
    else:
        out_dir = ROOT_DIR / CONFIG["paths"]["outputs"]
    out_dir.mkdir(parents=True, exist_ok=True)

    ts         = datetime.now().strftime("%Y%m%d_%H%M%S")
    safe_topic = re.sub(r'[^a-zA-Z0-9_\-]', '_', topic[:35]).strip("_")
    out_path   = out_dir / f"{safe_topic}_{theme.name.replace(' ','_')}_{ts}.pptx"

    prs.save(str(out_path))
    log.info(f"✔ Saved PPTX ({theme.name}): {out_path}")
    return str(out_path)
